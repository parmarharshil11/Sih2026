from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()

# Slide 1: Title Page
slide_layout_title = prs.slide_layouts[0]
slide1 = prs.slides.add_slide(slide_layout_title)
title1 = slide1.shapes.title
sub1 = slide1.placeholders[1]

title1.text = "SMART INDIA HACKATHON 2025"
title1.text_frame.paragraphs[0].font.size = Pt(36)
title1.text_frame.paragraphs[0].font.bold = True
title1.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

sub1.text = (
    "• Problem Statement ID: SIH-2025-CC-01\n"
    "• Problem Statement Title: CAPACITY CONNECT - A Digital Capacity Building and Learning Management Portal\n"
    "• Theme: Smart Education / Enterprise Solutions\n"
    "• PS Category: Software\n"
    "• Team ID: [Your Team ID]\n"
    "• Team Name: [Your Team Name]"
)
for p in sub1.text_frame.paragraphs:
    p.font.size = Pt(18)

def add_custom_slide(title_text, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(30)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    tf = slide.placeholders[1].text_frame
    tf.clear()
    
    for bullet, level, bold in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = level
        p.font.size = Pt(18)
        if bold:
            p.font.bold = True
    return slide

# Slide 2: Proposed Solution
s2_bullets = [
    ("CAPACITY CONNECT: Intelligent Capacity Building & Learning Portal", 0, True),
    ("Detailed Explanation:", 0, True),
    ("A centralized web portal that automates organizational competency tracking, matches employees with expert trainers, and delivers structured digital learning.", 1, False),
    ("How it Addresses the Problem:", 0, True),
    ("Replaces fragmented spreadsheets with real-time capability metrics, role benchmarks, and automated upskilling paths.", 1, False),
    ("Innovation & Uniqueness (USP):", 0, True),
    ("Automated Skill-Gap Engine: Evaluates (Required Level - Current Level) into Criticality tiers.", 1, False),
    ("6-Signal Weighted Matching Engine: Pairs trainees to trainers using skill overlap, availability, ratings & experience.", 1, False),
    ("Tamper-Proof QR Certification: Issues verifiable digital certificates with instant scan verification.", 1, False)
]
add_custom_slide("PROPOSED SOLUTION", s2_bullets)

# Slide 3: Technical Approach
s3_bullets = [
    ("Technologies & Frameworks:", 0, True),
    ("Frontend: Next.js 14 (App Router), TailwindCSS, Radix UI, Lucide Icons", 1, False),
    ("Backend API: NestJS (Enterprise Modular Architecture, RBAC & Guards)", 1, False),
    ("Database & Caching: PostgreSQL (Prisma ORM, 40+ models), Redis", 1, False),
    ("Storage & Infrastructure: MinIO / S3 Object Storage, Docker Containerization", 1, False),
    ("Implementation Process Flow:", 0, True),
    ("1. Trainee Self/Manager Assessment -> 2. Skill Gap Computation Engine ->", 1, False),
    ("3. Multi-Signal Trainer/Course Matching -> 4. Interactive Learning & Quiz -> 5. QR Certificate Verification", 1, False)
]
s3 = add_custom_slide("TECHNICAL APPROACH", s3_bullets)

# Slide 4: Feasibility and Viability
s4_bullets = [
    ("Feasibility Analysis:", 0, True),
    ("High Scalability: Microservice-ready Turborepo monorepo with containerized services.", 1, False),
    ("Enterprise Security: Argon2id password hashing, HTTP-only JWTs, Role-Based Access Control.", 1, False),
    ("Potential Challenges & Risks:", 0, True),
    ("Complex custom skill taxonomies across diverse corporate/academic departments.", 1, False),
    ("Scheduling conflicts between active trainers and workforce trainees.", 1, False),
    ("Mitigation Strategies:", 0, True),
    ("Dynamic Competency Framework: Flexible 40+ table schema allowing custom skill-to-competency mappings.", 1, False),
    ("Availability Slot Matching: Algorithm factors real-time weekly schedule matrices into pairing scores.", 1, False)
]
add_custom_slide("FEASIBILITY AND VIABILITY", s4_bullets)

# Slide 5: Impact and Benefits
s5_bullets = [
    ("Target Audience: Enterprises, Government Departments, Universities & Skill Centers.", 0, True),
    ("Key Benefits & Outcomes:", 0, True),
    ("Economic: Reduces external training expenditures by up to 60% through optimized internal peer training.", 1, False),
    ("Operational: Cuts manual HR administration and training scheduling time by 80%.", 1, False),
    ("Strategic / Management: Real-time Organizational Heatmaps spot skill deficits before critical project failures.", 1, False),
    ("Trainee Empowerment: Clear visibility into personal career growth and fraud-proof shareable credentials.", 1, False)
]
add_custom_slide("IMPACT AND BENEFITS", s5_bullets)

# Slide 6: Research and References
s6_bullets = [
    ("Industry Research & Validation:", 0, True),
    ("World Economic Forum (Future of Jobs Report): 50% of employees require significant reskilling by 2025.", 1, False),
    ("Gartner HR Insights: 70% of employees report lack of mastery in skills needed for future roles.", 1, False),
    ("Key References & Standards:", 0, True),
    ("National Skill Qualification Framework (NSQF) Competency Modeling Guidelines", 1, False),
    ("OWASP Application Security Verification Standard (ASVS) for RBAC & Token Security", 1, False)
]
add_custom_slide("RESEARCH AND REFERENCES", s6_bullets)

prs.save('Capacity_Connect_SIH2025.pptx')
print("Updated Capacity_Connect_SIH2025.pptx successfully!")
